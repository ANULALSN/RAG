from pathlib import Path

from pptx import Presentation


def clean_text(text: str) -> str:
    """Clean common PowerPoint extraction artifacts."""

    replacements = {
        "\u2642": "",
        "\u00a0": " ",
    }

    for old, new in replacements.items():
        text = text.replace(old, new)

    lines = []

    for line in text.splitlines():
        line = " ".join(line.split())

        if line:
            lines.append(line)

    return "\n".join(lines)


def extract_pptx(path: str | Path) -> list[dict]:
    """
    Extract textual content from every slide in a PPTX.

    Returns one normalized record per slide.
    """

    path = Path(path)

    if not path.exists():
        raise FileNotFoundError(
            f"PPTX not found: {path}"
        )

    if path.suffix.lower() != ".pptx":
        raise ValueError(
            f"Expected .pptx file, got: {path.suffix}"
        )

    presentation = Presentation(path)

    slides = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):
        text_parts = []
        title = ""

        for shape in slide.shapes:

            # Normal text boxes / placeholders
            if hasattr(shape, "text"):
                text = clean_text(shape.text)

                if not text:
                    continue

                # Detect PowerPoint title placeholder
                if (
                    hasattr(shape, "is_placeholder")
                    and shape.is_placeholder
                    and shape.placeholder_format.type == 1
                ):
                    if not title:
                        title = text

                    continue

                text_parts.append(text)

            # Extract table contents
            if shape.has_table:
                for row in shape.table.rows:
                    cells = []

                    for cell in row.cells:
                        cell_text = clean_text(cell.text)

                        if cell_text:
                            cells.append(cell_text)

                    if cells:
                        text_parts.append(
                            " | ".join(cells)
                        )

        slide_text = "\n".join(text_parts)

        # Remove duplicated title from body text.
        if title and slide_text:
            title_normalized = (
                " ".join(title.split()).casefold()
            )

            lines = slide_text.splitlines()

            if lines:
                first_line = (
                    " ".join(lines[0].split()).casefold()
                )

                if first_line == title_normalized:
                    slide_text = "\n".join(
                        lines[1:]
                    ).strip()

        # IMPORTANT:
        # This append MUST remain INSIDE the slide loop.
        slides.append(
            {
                "slide": slide_number,
                "title": title,
                "text": slide_text,
                "source_type": "pptx",
            }
        )

    return slides


if __name__ == "__main__":

    pptx_path = Path(
        "data/raw/BigData/Module1/Module 1_BD.pptx"
    )

    slides = extract_pptx(pptx_path)

    print("=" * 70)
    print("NORMALIZED PPTX EXTRACTION TEST")
    print("=" * 70)

    print(f"File: {pptx_path}")
    print(f"Slides: {len(slides)}")

    non_empty = sum(
        bool(
            slide.get("title")
            or slide.get("text")
        )
        for slide in slides
    )

    print(f"Non-empty slides: {non_empty}")
    print(
        f"Empty slides: "
        f"{len(slides) - non_empty}"
    )

    for slide in slides[:5]:
        print("\n" + "-" * 70)
        print(f"SLIDE {slide['slide']}")
        print(
            f"TITLE: "
            f"{slide['title'] or '[No title]'}"
        )
        print("-" * 70)
        print(
            slide["text"]
            or "[No textual content]"
        )